import os
import sys
from pathlib import Path
import tarfile
import shutil
import re
from xml.etree import ElementTree as et
from io import BytesIO
import json

import yaml

from invoke import task, call
import requests

with open('_config.yml', 'r', encoding='utf8') as yf:
    CONFIG = yaml.safe_load(yf)

BRANCH = 'archive'
CURRENT = CONFIG['sphinx']['config']['term']
BUILD_DIR = Path('_build')
FULL_DIR = BUILD_DIR / 'site'
CUR_DIR = BUILD_DIR / 'dirhtml'
CUR_DEPLOY_DIR = FULL_DIR / CURRENT

INVENTORY_CSV = 'https://docs.google.com/spreadsheets/d/e/2PACX-1vTfaz9ZdI_DNh3SpID1D71dCZrA1YIDraBlf9aou_hcqYLjudQd4Mw0GIpfL_ViYz73UfYhamtk9oEC/pub?gid=0&single=true&output=csv'


def _msg(fmt, *args):
    print(fmt % args, file=sys.stderr)


@task
def info(c):
    print('Current term:', CURRENT)


@task
def fetch_video_list(c, skip_if_exists=False):
    "Fetch list of videos"
    out = Path('videos/manifest.jsonl')
    id = os.environ['BUNNY_STREAMS_LIB']
    key = os.environ['BUNNY_STREAMS_KEY']
    params = {'itemsPerPage': 200}
    headers = {'AccessKey': key}
    res = requests.get(f'https://video.bunnycdn.com/library/{id}/videos', params=params, headers=headers)
    with out.open('w', encoding='utf8') as outf:
        for item in res.json()['items']:
            print(json.dumps(item), file=outf)


@task()
def build(c, rebuild=False, builder='dirhtml'):
    """
    Build the current site.
    """
    if rebuild:
        c.run(f'jb build --all --builder {builder} .')
    else:
        c.run(f'jb build --builder {builder} .')


@task
def serve(c, full_site=False, port=8000, rebuild=False):
    from livereload import Server

    def do_rebuild():
        _msg('building and copying site')
        build(c, rebuild=rebuild)
        if full_site:
            build_site(c)

    _ignore_re = re.compile(r'(^(\.[\\/])?_build|.*\..ipynb_checkpoints)[\\/]')

    def ignore(p):
        return _ignore_re.match(p)

    do_rebuild()
    srv_root = FULL_DIR if full_site else CUR_DIR
    server = Server()
    server.watch('.', do_rebuild, ignore=ignore)
    server.serve(port=port, root=srv_root)


@task
def copy_archive(c):
    archive_path = BUILD_DIR / 'archive.tar'
    p = archive_path.as_posix()

    FULL_DIR.mkdir(exist_ok=True, parents=True)

    _msg('creating git archive')
    c.run(f'git archive -o {p} archive')

    with tarfile.open(archive_path) as tf:
        _msg('extracting %s', archive_path)
        tf.extractall(FULL_DIR)


@task(pre=[copy_archive])
def copy_root(c):
    _msg('copying static root files')
    shutil.copytree('_root', FULL_DIR, dirs_exist_ok=True)


@task(pre=[build, copy_root])
def build_site(c):
    _msg('copying current tree to site target')
    shutil.copytree(CUR_DIR, CUR_DEPLOY_DIR, dirs_exist_ok=True)


@task
def clean(c):
    _msg('removing build directory')
    shutil.rmtree(BUILD_DIR)


@task
def extract_slides(c, dir):
    out_dir = Path('videos')
    if dir.endswith('"'):
        dir = dir[:-1]
    dir = Path(dir)
    import pptx
    _msg('scanning slides in %s', dir)
    for file in dir.glob("*.pptx"):
        print('reading', file.name)
        name = file.stem
        out_file = out_dir / f'{name}.slides.txt'
        ppt = pptx.Presentation(file)
        with open(out_file, 'w', encoding='utf8') as otf:
            for slide in ppt.slides:
                for shape in slide.shapes:
                    text = getattr(shape, 'text', None)
                    if text:
                        print(text, file=otf)


@task
def get_videos(c, folder):
    url = f'https://boisestate.hosted.panopto.com/Panopto/Podcast/Podcast.ashx?courseid={folder}&type=mp4'
    id_re = re.compile(r'https://.*/Syndication/([a-z0-9-]+)\.mp4')
    vidroot = Path('videos')
    _msg('fetching videos for folder %s', folder)
    res = requests.get(url)
    rss = et.parse(BytesIO(res.content))
    chan = rss.find('channel')

    for item in chan.findall('item'):
        title = item.find('title').text
        guid = item.find('guid').text.strip()
        m = id_re.match(guid)
        id = m.group(1)
        _msg('found video %s with id %s', title, id)

        srturi = f'https://boisestate.hosted.panopto.com/Panopto/Pages/Transcription/GenerateSRT.ashx?id={id}&language=0'
        srt = requests.get(srturi)
        _msg('downloaded subtitles, saving')
        file = vidroot / f'{title}.srt'
        file.write_bytes(srt.text.encode('utf8'))


@task
def onedrive_login(c):
    import toml
    from msal import PublicClientApplication

    odtp = Path('onedrive.toml')
    cfg = toml.loads(odtp.read_text())
    client_cfg = cfg['client']
    tenant = client_cfg['tenant']
    auth = f'https://login.microsoftonline.com/{tenant}/'

    app = PublicClientApplication(client_cfg['id'], authority=auth)
    result = app.acquire_token_interactive(scopes=["files.readwrite"])
    print(result)
