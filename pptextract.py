"""
Extract PowerPoint to text.

Based on https://stackoverflow.com/a/47272482/1385039
"""
from os import write
import sys
from pathlib import Path
import csv

from tinytag import TinyTag
import pptx

out_dir = Path('videos')

for file in sys.argv[1:]:
    file = Path(file)
    print('reading', file.name)
    name = file.stem
    out_file = out_dir / f'{name}.slides.txt'
    ppt = pptx.Presentation(file)
    with open(out_file, 'w', encoding='utf8') as otf:
        for slide in ppt.slides:
            for shape in slide.shapes:
                text = getattr(shape, 'text', None)
                if text:
                    print(text, file=otf)
